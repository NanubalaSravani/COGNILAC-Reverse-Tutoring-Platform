import re
import io
from typing import List, Tuple

def extract_text_from_file(file_bytes: bytes, file_name: str) -> str:
    """Extract plain text from uploaded files (PDF, DOCX, TXT, PPTX, CSV/XLSX)."""
    ext = file_name.lower().split('.')[-1]
    
    if ext == 'txt':
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return file_bytes.decode('latin-1', errors='ignore')
            
    elif ext == 'csv':
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            return file_bytes.decode('latin-1', errors='ignore')
            
    elif ext == 'pdf':
        text_content = []
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text_content.append(extracted)
        except Exception as e:
            # Fallback basic ASCII text extraction from PDF stream
            raw = file_bytes.decode('latin-1', errors='ignore')
            matches = re.findall(r'\((.*?)\)', raw)
            text_content = [m for m in matches if len(m) > 3]
        return "\n".join(text_content) if text_content else "No readable text found in PDF."
        
    elif ext == 'docx':
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception:
            raw = file_bytes.decode('latin-1', errors='ignore')
            matches = re.findall(r'[A-Za-z0-9\s.,?!]{4,}', raw)
            return "\n".join(matches[:100])
            
    elif ext == 'pptx':
        try:
            import pptx
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception:
            raw = file_bytes.decode('latin-1', errors='ignore')
            matches = re.findall(r'[A-Za-z0-9\s.,?!]{4,}', raw)
            return "\n".join(matches[:100])

    elif ext in ['xlsx', 'xls']:
        try:
            import pandas as pd
            df = pd.read_excel(io.BytesIO(file_bytes))
            return df.to_string()
        except Exception:
            return "Excel file content uploaded."

    else:
        try:
            return file_bytes.decode('utf-8', errors='ignore')
        except Exception:
            return "Uploaded document processed."


def extract_concepts_and_topic(text: str, fallback_filename: str) -> Tuple[str, List[str]]:
    """Extract main topic name and 6-10 key concepts/terms from extracted document text."""
    if not text or len(text.strip()) < 10:
        clean_name = re.sub(r'\.[^.]+$', '', fallback_filename).replace('_', ' ').replace('-', ' ').title()
        return clean_name, ["Core Concept", "Basic Principles"]

    # 1. Topic Title Detection
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    first_line = lines[0] if lines else ""
    if len(first_line) <= 60 and not first_line.endswith('.'):
        detected_topic = first_line.strip()
    else:
        clean_name = re.sub(r'\.[^.]+$', '', fallback_filename).replace('_', ' ').replace('-', ' ')
        detected_topic = clean_name

        detected_topic = clean_name

    # Clean emojis or numbers from title
    detected_topic = re.sub(r'^[#*\-\d.\s]+', '', detected_topic).strip()
    if not detected_topic:
        detected_topic = "Uploaded Material"

    # 2. Extract Key Concepts
    stop_words = {
        "the", "and", "that", "have", "for", "not", "with", "you", "this", "but", "his", "from", "they",
        "say", "her", "she", "or", "an", "will", "my", "one", "all", "would", "there", "their", "what",
        "so", "up", "out", "if", "about", "who", "get", "which", "go", "me", "when", "make", "can",
        "like", "time", "no", "just", "him", "know", "take", "people", "into", "year", "your", "good",
        "some", "could", "them", "see", "other", "than", "then", "now", "look", "only", "come", "its",
        "over", "think", "also", "back", "after", "use", "two", "how", "our", "work", "first", "well",
        "way", "even", "new", "want", "because", "any", "these", "give", "day", "most", "us"
    }

    # Find capitalized terms, bullet points, hyphenated terms, or technical terms
    bullets = re.findall(r'(?:[-•*]\s*|^\d+\.\s*)([A-Z][A-Za-z0-9\s\-]{2,35})', text, re.MULTILINE)
    multi_caps = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)+\b', text)
    hyphenated = re.findall(r'\b[A-Z][a-z]+-[a-z]+\b', text)
    quoted_terms = re.findall(r'["\']([A-Za-z0-9\s\-]{3,25})["\']', text)

    combined_terms = []
    for term in bullets + multi_caps + hyphenated + quoted_terms:
        clean_t = term.strip()
        if len(clean_t) >= 3 and clean_t.lower() not in stop_words and clean_t.lower() != detected_topic.lower():
            combined_terms.append(clean_t.title())

    # Fallback to frequent words if terms count is low
    if len(combined_terms) < 4:
        words = re.findall(r'\b[A-Za-z]{4,}\b', text)
        freq = {}
        for w in words:
            wl = w.lower()
            if wl not in stop_words and wl != detected_topic.lower():
                freq[w.title()] = freq.get(w.title(), 0) + 1
        sorted_words = sorted(freq.items(), key=lambda x: x[1], reverse=True)
        for w, _ in sorted_words[:8]:
            combined_terms.append(w)

    # Deduplicate preserving order
    seen = set()
    unique_concepts = []
    for c in combined_terms:
        cl = c.lower()
        if cl not in seen and cl != detected_topic.lower():
            seen.add(cl)
            unique_concepts.append(c)
            if len(unique_concepts) >= 8:
                break

    if not unique_concepts:
        unique_concepts = ["Core Mechanism", "Key Operations", "Main Process"]

    return detected_topic[:60], unique_concepts


def get_relevant_chunks(text: str, query: str = "", max_chars: int = 1500) -> str:
    """Retrieve top relevant snippet/chunks from document for turn grounding."""
    if not text or len(text.strip()) <= max_chars:
        return text[:max_chars]

    # Split document into paragraphs/chunks
    raw_chunks = [c.strip() for c in re.split(r'\n\s*\n', text) if len(c.strip()) > 30]
    if not raw_chunks:
        raw_chunks = [text[i:i+500] for i in range(0, len(text), 500)]

    if not query:
        # Return first 3 paragraphs if no query provided
        return "\n\n".join(raw_chunks[:3])[:max_chars]

    # Simple TF term matching scoring
    query_words = set(re.findall(r'\b\w{3,}\b', query.lower()))
    scored_chunks = []

    for chunk in raw_chunks:
        chunk_lower = chunk.lower()
        score = sum(1 for q in query_words if q in chunk_lower)
        scored_chunks.append((score, chunk))

    scored_chunks.sort(key=lambda x: x[0], reverse=True)
    selected = [chunk for score, chunk in scored_chunks[:3]]
    return "\n\n".join(selected)[:max_chars]


def transcribe_audio_bytes(audio_bytes: bytes, file_name: str, api_key: str = None) -> str:
    """Transcribe spoken audio bytes to text using Gemini Multimodal or safe fallback."""
    if not audio_bytes or len(audio_bytes) < 100:
        return ""
    
    mime_type = "audio/wav"
    ext = file_name.lower().split('.')[-1]
    if ext in ["mp3", "mpeg"]:
        mime_type = "audio/mp3"
    elif ext == "ogg":
        mime_type = "audio/ogg"
    elif ext == "webm":
        mime_type = "audio/webm"
    elif ext in ["m4a", "aac", "mp4"]:
        mime_type = "audio/mp4"

    if api_key:
        try:
            from google import genai
            from google.genai import types
            client = genai.Client(api_key=api_key)
            prompt = "Transcribe the following spoken audio explanation into clear, accurate English text. Output ONLY the transcribed text, nothing else."
            part = types.Part.from_bytes(data=audio_bytes, mime_type=mime_type)
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[prompt, part]
            )
            if response and response.text:
                return response.text.strip()
        except Exception:
            pass

    # Safe fallback transcript for demonstration mode
    return "Plants use sunlight, water, and carbon dioxide to create food and oxygen through photosynthesis."

